"""
Document processing tools

Provides document parsing and extraction capabilities:
- PDF text extraction
- Word documents (docx)
- Excel spreadsheets
- PowerPoint presentations
- JSON/XML parsing
- ZIP archive extraction

Most dependencies are optional with graceful degradation.
"""

import os
import json
from pathlib import Path
from typing import Dict, Any, Optional, List
from ..core.logger import get_logger
from .base import BaseTool, ToolOutput

logger = get_logger(__name__)

# Optional dependencies with graceful fallback
try:
    import pypdf
    HAS_PYPDF = True
except ImportError:
    HAS_PYPDF = False
    logger.debug("pypdf not installed. PDF extraction unavailable.")

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    logger.debug("python-docx not installed. DOCX extraction unavailable.")

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False
    logger.debug("pandas not installed. Excel extraction unavailable.")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.debug("python-pptx not installed. PPTX extraction unavailable.")


class DocumentExtractTool(BaseTool):
    """Universal document content extraction tool"""
    
    # Supported file extensions and their handlers
    SUPPORTED_EXTENSIONS = {
        'pdf': '_extract_pdf',
        'docx': '_extract_docx',
        'doc': '_extract_docx',
        'xlsx': '_extract_excel',
        'xls': '_extract_excel',
        'csv': '_extract_csv',
        'pptx': '_extract_pptx',
        'json': '_extract_json',
        'xml': '_extract_xml',
        'txt': '_extract_text',
        'md': '_extract_text',
        'py': '_extract_text',
        'js': '_extract_text',
        'ts': '_extract_text',
        'html': '_extract_text',
        'css': '_extract_text',
        'yaml': '_extract_text',
        'yml': '_extract_text',
    }
    
    def __init__(self, safety_guard=None, cache_dir: str = "tmp/documents"):
        super().__init__(
            name="document_extract",
            description="Извлечение содержимого из документов (PDF, Word, Excel, PowerPoint, JSON, XML, текстовые файлы)",
            safety_guard=safety_guard
        )
        self.cache_dir = cache_dir
        os.makedirs(cache_dir, exist_ok=True)
    
    def _get_file_extension(self, path: str) -> str:
        """Get lowercase file extension without dot"""
        return Path(path).suffix.lower().strip('.')
    
    async def execute(self, input_data: Dict[str, Any]) -> ToolOutput:
        """Extract content from document"""
        file_path = input_data.get("file_path")
        max_chars = input_data.get("max_chars", 50000)  # Limit to prevent token overflow
        
        if not file_path:
            return ToolOutput(success=False, result=None, error="file_path required")
        
        # Safety check
        if self.safety_guard:
            if not self.safety_guard.validate_path(file_path):
                return ToolOutput(success=False, result=None, error="Invalid or unsafe path")
        
        # Check if file exists
        path = Path(file_path)
        if not path.exists():
            return ToolOutput(success=False, result=None, error=f"File not found: {file_path}")
        
        # Get extension and handler
        ext = self._get_file_extension(file_path)
        
        if ext not in self.SUPPORTED_EXTENSIONS:
            supported = ", ".join(sorted(self.SUPPORTED_EXTENSIONS.keys()))
            return ToolOutput(
                success=False,
                result=None,
                error=f"Unsupported format: .{ext}. Supported: {supported}"
            )
        
        try:
            handler_name = self.SUPPORTED_EXTENSIONS[ext]
            handler = getattr(self, handler_name)
            content = await handler(file_path)
            
            # Truncate if needed
            original_length = len(content)
            if original_length > max_chars:
                content = content[:max_chars] + f"\n\n... [Truncated: {original_length} chars total, showing first {max_chars}]"
            
            return ToolOutput(
                success=True,
                result={
                    "file_path": str(path.absolute()),
                    "file_name": path.name,
                    "extension": ext,
                    "content": content,
                    "char_count": original_length,
                    "truncated": original_length > max_chars
                }
            )
            
        except Exception as e:
            logger.error(f"DocumentExtractTool error for {file_path}: {e}")
            return ToolOutput(success=False, result=None, error=str(e))
    
    async def _extract_pdf(self, path: str) -> str:
        """Extract text from PDF"""
        if not HAS_PYPDF:
            raise ImportError("pypdf not installed. Run: pip install pypdf")
        
        reader = pypdf.PdfReader(path)
        text_parts = []
        
        for i, page in enumerate(reader.pages, 1):
            page_text = page.extract_text()
            if page_text.strip():
                text_parts.append(f"--- Page {i} ---\n{page_text}")
        
        return "\n\n".join(text_parts)
    
    async def _extract_docx(self, path: str) -> str:
        """Extract text from Word document"""
        if not HAS_DOCX:
            raise ImportError("python-docx not installed. Run: pip install python-docx")
        
        doc = DocxDocument(path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
    
    async def _extract_excel(self, path: str) -> str:
        """Extract content from Excel file"""
        if not HAS_PANDAS:
            raise ImportError("pandas not installed. Run: pip install pandas openpyxl")
        
        xl = pd.ExcelFile(path)
        result_parts = []
        
        for sheet_name in xl.sheet_names:
            df = pd.read_excel(xl, sheet_name=sheet_name)
            result_parts.append(f"### Sheet: {sheet_name}")
            result_parts.append(f"Rows: {len(df)}, Columns: {len(df.columns)}")
            result_parts.append(f"Columns: {', '.join(df.columns.astype(str))}")
            result_parts.append("")
            
            # Show first 50 rows as markdown table
            if len(df) > 50:
                result_parts.append(df.head(50).to_markdown(index=False))
                result_parts.append(f"\n... [{len(df) - 50} more rows]")
            else:
                result_parts.append(df.to_markdown(index=False))
            
            result_parts.append("")
        
        return "\n".join(result_parts)
    
    async def _extract_csv(self, path: str) -> str:
        """Extract content from CSV file"""
        if not HAS_PANDAS:
            raise ImportError("pandas not installed. Run: pip install pandas")
        
        df = pd.read_csv(path)
        result_parts = [
            f"Rows: {len(df)}, Columns: {len(df.columns)}",
            f"Columns: {', '.join(df.columns.astype(str))}",
            ""
        ]
        
        # Show first 50 rows
        if len(df) > 50:
            result_parts.append(df.head(50).to_markdown(index=False))
            result_parts.append(f"\n... [{len(df) - 50} more rows]")
        else:
            result_parts.append(df.to_markdown(index=False))
        
        return "\n".join(result_parts)
    
    async def _extract_pptx(self, path: str) -> str:
        """Extract text from PowerPoint"""
        if not HAS_PPTX:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")
        
        prs = Presentation(path)
        text_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            
            if slide_texts:
                text_parts.append(f"--- Slide {i} ---")
                text_parts.extend(slide_texts)
                text_parts.append("")
        
        return "\n".join(text_parts)
    
    async def _extract_json(self, path: str) -> str:
        """Extract content from JSON file"""
        with open(path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, ensure_ascii=False, indent=2)
    
    async def _extract_xml(self, path: str) -> str:
        """Extract content from XML file"""
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Try to parse and format
        try:
            import xml.dom.minidom
            dom = xml.dom.minidom.parseString(content)
            return dom.toprettyxml(indent="  ")
        except Exception:
            # Return raw if parsing fails
            return content
    
    async def _extract_text(self, path: str) -> str:
        """Extract content from text-based files"""
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()


class ZipExtractTool(BaseTool):
    """Tool for extracting ZIP archives"""
    
    def __init__(self, safety_guard=None, extract_dir: str = "tmp/extracted"):
        super().__init__(
            name="zip_extract",
            description="Распаковка ZIP архивов",
            safety_guard=safety_guard
        )
        self.extract_dir = extract_dir
        os.makedirs(extract_dir, exist_ok=True)
    
    async def execute(self, input_data: Dict[str, Any]) -> ToolOutput:
        """Extract ZIP archive"""
        import zipfile
        
        file_path = input_data.get("file_path")
        target_dir = input_data.get("target_dir")
        
        if not file_path:
            return ToolOutput(success=False, result=None, error="file_path required")
        
        # Safety check
        if self.safety_guard:
            if not self.safety_guard.validate_path(file_path):
                return ToolOutput(success=False, result=None, error="Invalid or unsafe path")
        
        path = Path(file_path)
        if not path.exists():
            return ToolOutput(success=False, result=None, error=f"File not found: {file_path}")
        
        if not path.suffix.lower() == '.zip':
            return ToolOutput(success=False, result=None, error="Only .zip files are supported")
        
        try:
            # Determine extraction path
            if target_dir:
                extract_path = Path(target_dir)
            else:
                extract_path = Path(self.extract_dir) / path.stem
            
            extract_path.mkdir(parents=True, exist_ok=True)
            
            # Extract
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(extract_path)
            
            # List extracted files
            extracted_files = []
            for root, _, files in os.walk(extract_path):
                for file in files:
                    rel_path = os.path.relpath(os.path.join(root, file), extract_path)
                    extracted_files.append(rel_path)
            
            return ToolOutput(
                success=True,
                result={
                    "archive": str(path.absolute()),
                    "extract_path": str(extract_path.absolute()),
                    "file_count": len(extracted_files),
                    "files": extracted_files[:100],  # Limit to first 100
                    "truncated": len(extracted_files) > 100
                }
            )
            
        except zipfile.BadZipFile:
            return ToolOutput(success=False, result=None, error="Invalid or corrupted ZIP file")
        except Exception as e:
            logger.error(f"ZipExtractTool error: {e}")
            return ToolOutput(success=False, result=None, error=str(e))


class DocumentInfoTool(BaseTool):
    """Tool for getting document metadata"""
    
    def __init__(self, safety_guard=None):
        super().__init__(
            name="document_info",
            description="Получение информации о документе (размер, тип, метаданные)",
            safety_guard=safety_guard
        )
    
    async def execute(self, input_data: Dict[str, Any]) -> ToolOutput:
        """Get document information"""
        file_path = input_data.get("file_path")
        
        if not file_path:
            return ToolOutput(success=False, result=None, error="file_path required")
        
        # Safety check
        if self.safety_guard:
            if not self.safety_guard.validate_path(file_path):
                return ToolOutput(success=False, result=None, error="Invalid or unsafe path")
        
        path = Path(file_path)
        if not path.exists():
            return ToolOutput(success=False, result=None, error=f"File not found: {file_path}")
        
        try:
            stat = path.stat()
            ext = path.suffix.lower().strip('.')
            
            result = {
                "file_path": str(path.absolute()),
                "file_name": path.name,
                "extension": ext,
                "size_bytes": stat.st_size,
                "size_human": self._human_readable_size(stat.st_size),
                "modified": stat.st_mtime,
                "is_supported": ext in DocumentExtractTool.SUPPORTED_EXTENSIONS,
            }
            
            # Get additional metadata for specific types
            if ext == 'pdf' and HAS_PYPDF:
                try:
                    reader = pypdf.PdfReader(file_path)
                    result["page_count"] = len(reader.pages)
                    if reader.metadata:
                        result["metadata"] = {
                            "title": reader.metadata.title,
                            "author": reader.metadata.author,
                            "subject": reader.metadata.subject,
                        }
                except Exception:
                    pass
            
            elif ext in ('xlsx', 'xls') and HAS_PANDAS:
                try:
                    xl = pd.ExcelFile(file_path)
                    result["sheet_count"] = len(xl.sheet_names)
                    result["sheets"] = xl.sheet_names
                except Exception:
                    pass
            
            return ToolOutput(success=True, result=result)
            
        except Exception as e:
            logger.error(f"DocumentInfoTool error: {e}")
            return ToolOutput(success=False, result=None, error=str(e))
    
    def _human_readable_size(self, size: int) -> str:
        """Convert bytes to human readable format"""
        for unit in ['B', 'KB', 'MB', 'GB']:
            if size < 1024.0:
                return f"{size:.1f} {unit}"
            size /= 1024.0
        return f"{size:.1f} TB"

